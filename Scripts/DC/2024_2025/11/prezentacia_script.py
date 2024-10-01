from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

# Create a PowerPoint presentation
prs = Presentation()

# Slide 1: Title Slide
slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]
title.text = "Prezentácia na tému 3D tlačiarne pre ZŠ na Hôrke"
subtitle.text = "Rozsah: 45 minút\nPočet žiakov: do 15\nTlačiareň: Prusa i3 MK3"

# Slide 2: Úvod
slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
title, content = slide_2.shapes.title, slide_2.placeholders[1]
title.text = "1. Úvod (10 minút)"
content.text = ("Krátky prehľad o 3D tlačiarniach, ich význame a základných princípoch.\n\n"
                "Vysvetlenie procesu 3D tlače a jeho rozdielov od tradičných výrobných procesov.\n\n"
                "Diskusia o rýchlom prototypovaní a potenciáli pre prispôsobené dizajny.\n\n"
                "Budúci potenciál 3D tlače vrátane domáceho použitia a zmeny výrobných procesov.\n\n"
                "Diskusia a návrhy od žiakov o budúcom využití 3D tlače.")

# Slide 3: Základy 3D Modelovania
slide_3 = prs.slides.add_slide(prs.slide_layouts[1])
title, content = slide_3.shapes.title, slide_3.placeholders[1]
title.text = "2. Základy 3D Modelovania (15 minút)"
content.text = ("Ukážka základov 3D modelovania na projektore.\n\n"
                "Diskusia o premenené návrhu modelu na tlačený objekt.\n\n"
                "Predstavenie softvérov Tinkercad a 3D Skicár pre jednoduché 3D modelovanie.")

# Slide 4: Interaktívna Skupinová Činnosť
slide_4 = prs.slides.add_slide(prs.slide_layouts[1])
title, content = slide_4.shapes.title, slide_4.placeholders[1]
title.text = "3. Interaktívna Skupinová Činnosť (15 minút)"
content.text = ("Žiaci vytvoria vlastné návrhy modelov na počítači.\n\n"
                "Téma pre návrhy môže byť zadaná alebo môžu mať žiaci voľný výber.")

# Slide 5: Diskusia a Otázky
slide_5 = prs.slides.add_slide(prs.slide_layouts[1])
title, content = slide_5.shapes.title, slide_5.placeholders[1]
title.text = "4. Diskusia a Otázky (5 minút)"
content.text = ("Priestor na otázky a odpovede o 3D tlači.\n\n"
                "Povzbudenie žiakov k pýtaní sa na proces tlače a možnosti 3D tlače.")

# Slide 6: Záver
slide_6 = prs.slides.add_slide(prs.slide_layouts[1])
title, content = slide_6.shapes.title, slide_6.placeholders[1]
title.text = "Záver (na konci dňa okolo 13:00)"
content.text = ("Vyberieme jeden návrh od žiakov a vytlačíme ich na 3D tlačiarni Prusa i3 MK3.")

# Save the presentation
pptx_fp = 'Prezentacia_3D_Tlaciarne_ZS_Horka.pptx'
prs.save(pptx_fp)

pptx_fp

